import os, docx
from pptx import Presentation
from pptx.shapes.autoshape import AutoShapeType
import pptx.shapes.autoshape
from docx.shared import Pt
from glob import glob

text_runs = []
para_runs = []
erg = 4
def getText(): ##gets text and stores in text_runs[]
    filename = glob("C:/Users/Armne/OneDrive/Կարմիր Երգարան/" + str(erg) + "*.pptx")[0] #gets name C:\Users\Armne\OneDrive\Կարմիր Երգարան
    prs = Presentation(filename)
    for slide in prs.slides:
        for shape in slide.shapes:
            # for x in AutoShapeType.mro():
            print( AutoShapeType(prs))
            if not (shape.has_text_frame):
                continue
            for paragraph in shape.text_frame.paragraphs:
                para_runs.append(paragraph.text)
                for run in paragraph.runs:
                    text_runs.append(run.text)


getText()

for i in text_runs:
    print(i)
for x in para_runs:
    print(i)
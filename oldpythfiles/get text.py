from pptx import Presentation
import time
prs = Presentation("C:/Users/moses/OneDrive/power point songs/28.pptx")
# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []
para_runs = []
def getText():
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_runs.append(paragraph.text)
                    for run in paragraph.runs:
                        text_runs.append(run.text)


##for i in para_runs:
##    print(i)
##    time.sleep(1)

##def translateAM(cur_text):
##    Arial = ""
##    ##Goto first box
##    sbox = driver.find_element_by_id("convertr_box") # Finds first box
##    sleep(0.5)
##    sbox.clear()
##    sbox.send_keys(cur_text) # the searchterm, uses "searchterm" var
##    sleep(1)
##    sbox.click()
##
##    #Goto Second box
##    obox = driver.find_element_by_id("unicode_convertr_box")
##    obox.click()
##    Arial = obox.text
##    print(Arial)
##    return Arial

def searchReplace(output):
    """"search and replace text in PowerPoint while preserving formatting"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_runs.append(paragraph.text)
                    for run in paragraph.runs:              
                        cur_text = run.text
##                        text_frame = shape.text_frame
##                        cur_text = text_frame.paragraphs[0].runs[0].text #.text fromhttps://python-pptx.readthedocs.io/en/latest/api/shapes.html#shape-objects-autoshapes
##--------------This part does the actual replacing part, not vital in finding the lines--------------
                        new_text = translateAM(cur_text)
                        text_frame.paragraphs[0].runs[0].text = new_text

    prs.save(output)

erg = "27" ## use to iterate through entire thang
prs = Presentation("C:/Users/moses/OneDrive/power point songs/" + erg +".pptx")

searchReplace("C:/Users/moses/OneDrive/power point songs/" + erg + " CONVERTED!" +".pptx")

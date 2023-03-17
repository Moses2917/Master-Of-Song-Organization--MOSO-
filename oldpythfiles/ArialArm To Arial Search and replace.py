##import Selenium
import time
from time import sleep
from selenium import webdriver
import pyautogui
from pyautogui import hotkey
from webdriver_manager.chrome import ChromeDriverManager
from pptx import Presentation
from selenium import webdriver
from webdriver_manager.chrome import ChromeDriverManager
from pptx.util import Inches, Pt
import os.path
## tip use pip install python-pptx to install the pptx module

option = webdriver.ChromeOptions()
option.add_argument('headless')
driver = webdriver.Chrome(ChromeDriverManager().install(),options=option)


#txt = open('C:\Users\Armnewlifechurch\Documents\Code\Python\67.txt', encoding='utf-8')
#tun = txt.read()
#print(tun)

##lines = []
##with open('67.txt', encoding='utf8') as txt:
##    for line in txt:
##        lines.append(line)
##print(lines)

##output of it ['1. öáñÓ³ÝùÝ»ñáõ Ù»ç ³ÝÑ³Ù³ñ, \x0bºí ÷áõß»ñáõ Ù»ç ëñ³Í³Ûñ, \x0bÊáñÑáõñ¹ Ù°³Ýáõß Ï³ÝóÝÇ Ùïù»ë, \x0bÂ» ¸áõÝ ½Çë ÏÑÇß»ë: \n']




webpage = r"http://www.armdict.com/tool/unicode-converter/" # edit me
##driver = webdriver.Chrome()
driver.get(webpage)


def translateAM(cur_text):
    Arial = ""
    ##Goto first box
    sbox = driver.find_element_by_id("convertr_box") # Finds first box
    sleep(0.5)
    sbox.clear()
    sbox.send_keys(cur_text) # the searchterm, uses "searchterm" var
    sleep(1)
    sbox.click()

    #Goto Second box
    obox = driver.find_element_by_id("unicode_convertr_box")
    obox.click()
    Arial = obox.text
##    print(Arial)
    return Arial
    
    #Get text, in clipboard
    #sleep(0.5)
    #pyautogui.hotkey('ctrl', 'a')
    #pyautogui.hotkey('ctrl', 'c')
    #sleep(0.5)
    #Arial = obox.text
    #print(Arial)

##def searchReplace(output):
##    """"search and replace text in PowerPoint while preserving formatting"""
##    #Useful Links ;)
##    #https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
##    #https://stackoverflow.com/questions/45247042/how-to-keep-original-text-formatting-of-text-with-python-powerpoint
##https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
##    for slide in prs.slides:
##        for shape in slide.shapes:
##            if shape.has_text_frame:
##                text_frame = shape.text_frame
##                cur_text = text_frame.paragraphs[0].runs[0].text
####------------- This part does the actual replacing part, not vital in finding the lines-------
##                new_text = translateAM(cur_text)
##                text_frame.paragraphs[0].runs[0].text = new_text
##
##    prs.save(output)

def searchReplace(output):
    """"search and replace text in PowerPoint while preserving formatting"""
##    new_text = ""
    for slide in prs.slides:
##        print("Starting again")
        for shape in slide.shapes:
##            print("Starting againss")
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
##                    print("Starting agin")
##                    print(new_text)
                    for run in paragraph.runs:
                        text_frame = shape.text_frame
                        cur_text = run.text
##                        text_frame = shape.text_frame
##                        cur_text = text_frame.paragraphs[0].runs[0].text #.text fromhttps://python-pptx.readthedocs.io/en/latest/api/shapes.html#shape-objects-autoshapes
##--------------This part does the actual replacing part, not vital in finding the lines--------------
                        new_text = cur_text.replace(cur_text, translateAM(cur_text))
##                        text_frame.paragraphs[0].runs[0].text = new_text
                        text_frame.add_paragraph().text = new_text
                        print(new_text)

    prs.save(output)



 

text_runs = []
para_runs = []
def getText(): ##gets text and stores in text_runs[]
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not (shape.has_text_frame):
                continue
            for paragraph in shape.text_frame.paragraphs:
                para_runs.append(paragraph.text)
                for run in paragraph.runs:
                    text_runs.append(run.text)
                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = translateAM(cur_text)
                    text_frame.paragraphs[0].runs[0].text = new_text


##ArialArm = u"1. ØÇ ÃáÕáõñ ½Çë, áñ ²ëïí³Íë"
#searchterm = "1. ØÇ ÃáÕáõñ ½Çë, áñ ²ëïí³Í" # edit me

def bruteForce():
    first = input("First slide: ")
    chorus = input("Chorus: ")
    second  = input("Second slide: ")
    third  = input("Third slide: ")
    forth  = input("Forth slide: ")
    fifth  = input("Fifth slide: ")
    sixth  = input("Sixth slide: ")
    sevent  = input("Sevent slide: ")
    eight  = input("Eight slide: ")
    nineth  = input("Nineth slide: ")
    tenth  = input("Tenth slide: ")



    slides = [first, chorus, second, third, forth, fifth, sixth, sevent, eight, nineth, tenth]
    arial = [first, chorus, second, third, forth, fifth, sixth, sevent, eight, nineth, tenth]

    for i in slides:
        searchterm = i
        if searchterm != "":
            translateAM()
        elif searchterm == "":
            if chorus != "":
                break
fileNum = 26 ## 26 bc 1-25 already done!!!
while True:
    yerg = str(fileNum) ## use to iterate through entire thang
    filePath = "C:/Users/moses/OneDrive/power point songs/" + yerg +".pptx"
    print(filePath)
    if fileNum > 901:
        break
    else:
        if os.path.exists(filePath):
            prs = Presentation(filePath)
            searchReplace("C:/Users/moses/OneDrive/Translated Songs/" + yerg +".pptx")
        else:
            print("No File Found")
    fileNum += 1
##  keep on going until fileNum reaches 901


##if os.path.exists("C:/Users/moses/OneDrive/power point songs/" + erg +".pptx"):
##   print("Exists")
##else:
##   print("Nay")


##for i in text_runs:
##    print(i)

##getText()
####engine to my car
##for i in text_runs:
##    searchterm = i
##    if searchterm != "":
##        translateAM()
##    elif searchterm == "":
##        if chorus != "":
##            break




##Two ways, either output here and copy paste later(prob), or store in a database of sorts

##begin process


##End the session



##driver.key_down(obox, control).send_keys('c').key_up(Keys.CONTROL).perform





####Goto first box
##sbox = driver.find_element_by_id("convertr_box")
##sleep(0.5)
##sbox.clear()
##sbox.send_keys(searchterm)
##sleep(1)
##sbox.click()
##
###Goto Second box
##obox = driver.find_element_by_id("unicode_convertr_box")
##obox.click()
##
###Get text
##sleep(0.5)
##pyautogui.hotkey('ctrl', 'a')
##pyautogui.hotkey('ctrl', 'c')
##sleep(0.5)
##Arial = obox.text
##print(Arial)



driver.quit()








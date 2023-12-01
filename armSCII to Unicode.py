def armSCII_To_Unicode(text):
    #translated into python by Movses Movsesyan, code from armdict.com 
    fromLettersStr = "² ´ ¶ ¸ º ¼ ¾ À Â Ä Æ È Ê Ì Î Ð Ò Ô Ö Ø Ú Ü Þ à â ä æ è ê ì î ð ò ô ö ø ¨ ú ü ³ µ • ¹ » ½ ¿ Á Ã Å Ç É Ë Í Ï Ñ Ó Õ × Ù Û Ý ß á ã å ç é ë í ï ñ ó õ ÷ ù ¨ û ý"
    fromLettersStr2 = "Ա Բ Գ Դ Ե Զ Է Ը Թ Ժ Ի Լ Խ Ծ Կ Հ Ձ Ղ Ճ Մ Յ Ն Շ Ո Չ Պ Ջ Ռ Ս Վ Տ Ր Ց Ւ Փ Ք և Օ Ֆ ա բ գ դ ե զ է ը թ ժ ի լ խ ծ կ հ ձ ղ ճ մ յ ն շ ո չ պ ջ ռ ս վ տ ր ց ւ փ ք և օ ֆ"
    toLettersStr = "Ա Բ Գ Դ Ե Զ Է Ը Թ Ժ Ի Լ Խ Ծ Կ Հ Ձ Ղ Ճ Մ Յ Ն Շ Ո Չ Պ Ջ Ռ Ս Վ Տ Ր Ց Ւ Փ Ք և Օ Ֆ ա բ գ դ ե զ է ը թ ժ ի լ խ ծ կ հ ձ ղ ճ մ յ ն շ ո չ պ ջ ռ ս վ տ ր ց ւ փ ք և օ ֆ"
    toLettersStr2 = "² ´ ¶ ¸ º ¼ ¾ À Â Ä Æ È Ê Ì Î Ð Ò Ô Ö Ø Ú Ü Þ à â ä æ è ê ì î ð ò ô ö ø ¨ ú ü ³ µ · ¹ » ½ ¿ Á Ã Å Ç É Ë Í Ï Ñ Ó Õ × Ù Û Ý ß á ã å ç é ë í ï ñ ó õ ÷ ù ¨ û ý"
    # text = '2. Ðñ»ßï³ÏÝ Çç³í §ÐÇëáõë¦,- Áë³í, §¸Ý»ë ³Ýáñ ³ÝáõÝÁ, ø³Ý½Ç Ù³Ñ»Ý åÇïÇ ÷ñÏ», ´áÉáñ Çñ ÅáÕáíáõñ¹Á¦: '#example string in armSCII


    # to unicode
    fromLetters = fromLettersStr.split(" ")
    fromLetters2 = fromLettersStr2.split(" ")
    for index in range(len(fromLetters)):
        text = text.replace(fromLetters[index], fromLetters2[index])
    # ### #
    text = text.replace("°", "՛")
    text = text.replace("¯", "՜")
    text = text.replace("ª", "՝")
    text = text.replace("±", "՞")
    text = text.replace("£", "։")
    # separate
    text = text.replace("«", ",")
    text = text.replace("§", "«")
    text = text.replace("¦", "»")
    text = text.replace("®", "…")
    text = text.replace("©", ".")
    #
    text = text.replace("·", "•")
    text = text.replace("¹", "»")
    text = text.replace("½", "¿")
    text = text.replace("¿", "Á")
    text = text.replace("Ã", "Å")
    text = text.replace("Ç", "É")
    text = text.replace("Ë", "Í")
    text = text.replace("Ï", "Ñ")
    text = text.replace("Ó", "×")
    text = text.replace("Ù", "Û")
    text = text.replace("Ý", "ß")
    text = text.replace("á", "ã")
    text = text.replace("å", "ç")
    text = text.replace("é", "ë")
    text = text.replace("í", "ï")
    text = text.replace("ñ", "ó")
    text = text.replace("õ", "÷")
    text = text.replace("ù", "û")
    text = text.replace("ý", "ÿ")
    text = text.replace("","\n")

    return text

from pptx import Presentation

def getText(erg): ##gets text and stores in text_runs[]
    prs = Presentation(r'C:\Users\{}\OneDrive\power point songs/{}'.format(user,erg))
    
    song = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not (shape.has_text_frame):
                continue
            for paragraph in shape.text_frame.paragraphs:
                # print(armSCII_To_Unicode(paragraph.text)+"\n")
                
                song.append(armSCII_To_Unicode(paragraph.text)+"\n")
                # for run in paragraph.runs:
                #     print(run.text)
    return song

from glob import glob
from os import environ
import docx
from docx.shared import Pt
from re import sub

user = environ.get("USERNAME")
Powerpoints=glob("*.pptx",root_dir=r'C:\Users\{}\OneDrive\power point songs/'.format(user))
for x in Powerpoints:
    full_Song = getText(x)
    fileName = sub("pptx?$","docx",x)
    
    doc = docx.Document()
    # doc.add_paragraph(full_Song)
    for paragraph in full_Song:
        doc.add_paragraph(paragraph)
    font = doc.styles['Normal'].font
    font.name = 'Arial'
    font.size = Pt(22)
    doc.save(r'C:\Users\{}\OneDrive\Word songs\pptSong/{}'.format(user, fileName))

